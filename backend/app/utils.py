import os
import pdfplumber
import docx
from pptx import Presentation
from pathlib import Path

def save_upload_file(upload_file, dest_folder="uploads"):
    os.makedirs(dest_folder, exist_ok=True)
    file_path = os.path.join(dest_folder, upload_file.filename)
    with open(file_path, "wb") as buffer:
        buffer.write(upload_file.file.read())
    return file_path

def extract_text(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(file_path)
    elif ext == ".docx":
        return extract_text_docx(file_path)
    elif ext == ".pptx":
        return extract_text_pptx(file_path)
    else:
        raise ValueError("Unsupported file type")

def extract_text_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

def extract_text_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_text_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def chunk_text(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = words[i:i+chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap
    return chunks 